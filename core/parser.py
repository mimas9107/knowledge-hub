"""
Knowledge Hub 文件解析模組

支援格式：PDF、PPTX、DOCX、Markdown
"""
from pathlib import Path
from typing import Optional, Dict, List

def parse_document(filepath: str) -> Optional[Dict]:
    """
    解析文件，提取文字內容
    
    Args:
        filepath: 文件路徑
    
    Returns:
        dict: {
            'text': 完整文字內容,
            'pages': 頁數（如適用）,
            'metadata': 文件元資料
        }
    """
    filepath = Path(filepath)
    
    if not filepath.exists():
        return None
    
    ext = filepath.suffix.lower()
    
    parsers = {
        '.pdf': parse_pdf,
        '.pptx': parse_pptx,
        '.docx': parse_docx,
        '.md': parse_markdown,
    }
    
    parser = parsers.get(ext)
    if parser:
        try:
            return parser(filepath)
        except Exception as e:
            print(f"解析失敗 {filepath}: {e}")
            return None
    
    return None


def parse_pdf(filepath: Path) -> Dict:
    """解析 PDF 文件"""
    try:
        import pdfplumber
    except ImportError:
        raise ImportError("請安裝 pdfplumber: pip install pdfplumber")
    
    text_parts = []
    metadata = {}
    
    with pdfplumber.open(filepath) as pdf:
        metadata['pages'] = len(pdf.pages)
        metadata['title'] = pdf.metadata.get('Title', '')
        metadata['author'] = pdf.metadata.get('Author', '')
        
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append({
                    'page': i + 1,
                    'text': page_text.strip()
                })
    
    return {
        'text': '\n\n'.join([p['text'] for p in text_parts]),
        'pages': text_parts,
        'metadata': metadata
    }


def parse_pptx(filepath: Path) -> Dict:
    """解析 PPTX 文件"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("請安裝 python-pptx: pip install python-pptx")
    
    prs = Presentation(filepath)
    text_parts = []
    
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_texts.append(shape.text.strip())
        
        if slide_texts:
            text_parts.append({
                'page': i + 1,
                'text': '\n'.join(slide_texts)
            })
    
    return {
        'text': '\n\n'.join([p['text'] for p in text_parts]),
        'pages': text_parts,
        'metadata': {
            'pages': len(prs.slides),
            'title': '',
            'author': ''
        }
    }


def parse_docx(filepath: Path) -> Dict:
    """解析 DOCX 文件"""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("請安裝 python-docx: pip install python-docx")
    
    doc = Document(filepath)
    
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())
    
    # DOCX 沒有頁碼概念，整份視為一頁
    return {
        'text': '\n\n'.join(paragraphs),
        'pages': [{'page': 1, 'text': '\n\n'.join(paragraphs)}],
        'metadata': {
            'pages': 1,
            'title': doc.core_properties.title or '',
            'author': doc.core_properties.author or ''
        }
    }


def parse_markdown(filepath: Path) -> Dict:
    """解析 Markdown 文件"""
    with open(filepath, 'r', encoding='utf-8') as f:
        text = f.read()
    
    # 嘗試從第一行 H1 標題取得標題
    title = ''
    lines = text.split('\n')
    for line in lines:
        if line.startswith('# '):
            title = line[2:].strip()
            break
    
    return {
        'text': text,
        'pages': [{'page': 1, 'text': text}],
        'metadata': {
            'pages': 1,
            'title': title,
            'author': ''
        }
    }
